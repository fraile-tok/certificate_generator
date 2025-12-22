import csv
from pathlib import Path
import subprocess

from pptx import Presentation


# CONFIG
TEMPLATE = Path("coupon.pptx")
CSV_FILE = Path("test.csv")

AUTO_PDF = False # Convert to PDF
AUTO_PNG = True # Convert to PNG

OUT_PPTX_DIR = Path("out_pptx")
OUT_PDF_DIR = Path("out_pdf")
OUT_PNG_DIR = Path("out_png")

OUT_PPTX_DIR.mkdir(exist_ok=True)


# FUCTIONS
def fill_name(template_path: Path, output_path: Path, full_name: str):
    prs = Presentation(str(template_path))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{{NAME}}" in run.text:
                        run.text = run.text.replace("{{NAME}}", full_name)

    prs.save(str(output_path))

def convert_to_pdf(pptx_path: Path, pdf_path: Path):
    OUT_PDF_DIR.mkdir(exist_ok=True)
    
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

    subprocess.run(
        [
            soffice_path,
            "--headless",
            "--convert-to",
            "pdf",
            str(pptx_path),
            "--outdir",
            str(pdf_path),
        ],
        check=True,
    )

    try:
        pptx_path.unlink()
        print(f"Deleted PPTX: {pptx_path}")
    except Exception as e:
        print(f"Could not delete {pptx_path}: {e}")

def convert_to_png(pptx_path: Path, jpg_dir: Path):
    OUT_PNG_DIR.mkdir(exist_ok=True)
    
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

    subprocess.run(
        [
            soffice_path,
            "--headless",
            '--convert-to', 'png',
            str(pptx_path),
            "--outdir",
            str(jpg_dir)
        ],
        check=True,
    )

    try:
        pptx_path.unlink()
        print(f"Deleted PPTX: {pptx_path}")
    except Exception as e:
        print(f"Could not delete {pptx_path}: {e}")

def main():
    with CSV_FILE.open(newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)

        for row in reader:
            # Merge first and last name
            first = row["NOMBRE"].strip()
            last = row["APELLIDO"].strip()
            full_name = f"{first} {last}".strip()

            file_name = (
                full_name.replace(" ", "_")
                .replace("/","-")
                .replace("\\", "-")
            )

            pptx_out = OUT_PPTX_DIR / f"{file_name}.pptx"

            print(f"Generating PPTX for {full_name}")
            fill_name(TEMPLATE, pptx_out, full_name)

            if AUTO_PDF:
                print(f"  -> Converting to PDF")
                convert_to_pdf(pptx_out, OUT_PDF_DIR)

            if AUTO_PNG:
                print(f"  -> Converting to PNG")
                convert_to_png(pptx_out, OUT_PNG_DIR)

if __name__ == "__main__":
    main()